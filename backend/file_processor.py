"""File processor for extracting text from documents and processing images."""

import base64
import io
from typing import Tuple, Optional
from pathlib import Path

# Supported file extensions
DOCUMENT_EXTENSIONS = {'.pdf', '.docx', '.pptx'}
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
SUPPORTED_EXTENSIONS = DOCUMENT_EXTENSIONS | IMAGE_EXTENSIONS


def get_file_type(filename: str) -> Optional[str]:
    """
    Get the type of file based on extension.
    Returns 'document', 'image', or None if unsupported.
    """
    ext = Path(filename).suffix.lower()
    if ext in DOCUMENT_EXTENSIONS:
        return 'document'
    elif ext in IMAGE_EXTENSIONS:
        return 'image'
    return None


def is_supported_file(filename: str) -> bool:
    """Check if the file type is supported."""
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def extract_from_pdf(file_content: bytes) -> str:
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader
    
    try:
        pdf_file = io.BytesIO(file_content)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        
        return "\n\n".join(text_parts) if text_parts else "No text could be extracted from this PDF."
    except Exception as e:
        return f"Error extracting PDF text: {str(e)}"


def extract_from_docx(file_content: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document
    
    try:
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(f"[Table Row] {row_text}")
        
        return "\n\n".join(paragraphs) if paragraphs else "No text could be extracted from this document."
    except Exception as e:
        return f"Error extracting DOCX text: {str(e)}"


def extract_from_pptx(file_content: bytes) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation
    
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
            
            if len(slide_parts) > 1:  # More than just the header
                slides_text.append("\n".join(slide_parts))
        
        return "\n\n".join(slides_text) if slides_text else "No text could be extracted from this presentation."
    except Exception as e:
        return f"Error extracting PPTX text: {str(e)}"


def process_image(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    Process an image file for vision model analysis.
    Returns (base64_data, mime_type)
    """
    from PIL import Image
    
    ext = Path(filename).suffix.lower()
    
    # Map extensions to MIME types
    mime_types = {
        '.png': 'image/png',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.gif': 'image/gif',
        '.bmp': 'image/bmp',
        '.webp': 'image/webp'
    }
    
    mime_type = mime_types.get(ext, 'image/png')
    
    # Optionally resize large images to reduce payload size
    try:
        img = Image.open(io.BytesIO(file_content))
        
        # Resize if too large (max 2048x2048)
        max_size = 2048
        if img.width > max_size or img.height > max_size:
            img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            
            # Convert back to bytes
            output = io.BytesIO()
            img_format = 'PNG' if ext == '.png' else 'JPEG'
            img.save(output, format=img_format)
            file_content = output.getvalue()
    except Exception as e:
        print(f"Warning: Could not resize image: {e}")
    
    base64_data = base64.b64encode(file_content).decode('utf-8')
    return base64_data, mime_type


def process_file(file_content: bytes, filename: str) -> dict:
    """
    Main entry point for processing files.
    
    Returns:
        dict with keys:
        - 'success': bool
        - 'type': 'document' or 'image'
        - 'text': extracted text (for documents)
        - 'image_data': base64 data (for images)
        - 'mime_type': MIME type (for images)
        - 'error': error message if failed
    """
    if not is_supported_file(filename):
        ext = Path(filename).suffix.lower()
        return {
            'success': False,
            'error': f"Unsupported file type: {ext}. Supported types: PDF, DOCX, PPTX, PNG, JPEG, GIF, BMP, WEBP"
        }
    
    file_type = get_file_type(filename)
    ext = Path(filename).suffix.lower()
    
    try:
        if file_type == 'document':
            if ext == '.pdf':
                text = extract_from_pdf(file_content)
            elif ext == '.docx':
                text = extract_from_docx(file_content)
            elif ext == '.pptx':
                text = extract_from_pptx(file_content)
            else:
                return {'success': False, 'error': f"Unknown document type: {ext}"}
            
            return {
                'success': True,
                'type': 'document',
                'text': text,
                'filename': filename
            }
        
        elif file_type == 'image':
            base64_data, mime_type = process_image(file_content, filename)
            return {
                'success': True,
                'type': 'image',
                'image_data': base64_data,
                'mime_type': mime_type,
                'filename': filename
            }
    
    except Exception as e:
        return {
            'success': False,
            'error': f"Error processing file: {str(e)}"
        }
